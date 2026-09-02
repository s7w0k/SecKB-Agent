"""Office/PDF renderers：docx / xlsx / pptx 由真实库结构校验生成。

python-docx / openpyxl / python-pptx / reportlab 生成真实二进制文件；
绝不只改标题。DOCX/PDF/PPTX/XLSX 抽样由 P3 做结构校验与视觉检查记录。
"""
from __future__ import annotations

from pathlib import Path

from scripts.enterprise_rag.renderers.markdown import blocks_to_markdown
from scripts.enterprise_rag.renderers.content import RenderedDoc


def render_docx(doc: RenderedDoc, out: Path) -> None:
    from docx import Document
    d = Document()
    d.add_heading(doc.title, 0)
    for b in doc.blocks:
        if b.kind == "heading":
            d.add_heading(b.text, 1)
        elif b.kind == "para":
            d.add_paragraph(b.text)
        elif b.kind in ("warning", "clause"):
            d.add_paragraph("> " + b.text)
        elif b.kind == "step":
            d.add_paragraph(b.text)
        elif b.kind == "table":
            if b.headers:
                t = d.add_table(rows=1, cols=len(b.headers))
                for j, h in enumerate(b.headers):
                    t.rows[0].cells[j].text = str(h)
                for r in b.rows:
                    cells = t.add_row().cells
                    for j, val in enumerate(r):
                        if j < len(cells):
                            cells[j].text = str(val)
        elif b.kind == "qa":
            d.add_paragraph("**Q**: " + b.question)
            d.add_paragraph("**A**: " + b.answer)
        elif b.kind == "log":
            for line in (b.items or []):
                d.add_paragraph(line)
    d.save(str(out))


def render_xlsx(doc: RenderedDoc, out: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = doc.family[:28] or "Sheet"
    ws.append(["section", "content"])
    for b in doc.blocks:
        if b.kind in ("heading", "para", "warning", "clause", "step", "prereq"):
            ws.append([b.kind, b.text])
        elif b.kind == "table":
            if b.headers:
                ws.append(["table_header"] + list(b.headers))
            for r in b.rows:
                ws.append(["row"] + [str(x) for x in r])
        elif b.kind == "qa":
            ws.append(["question", b.question])
            ws.append(["answer", b.answer])
    wb.save(str(out))


def render_pptx(doc: RenderedDoc, out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = doc.title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    first = True
    for b in doc.blocks:
        if b.kind in ("heading", "para", "clause", "warning", "step"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = b.text
            run.font.size = Pt(16 if b.kind == "heading" else 12)
        elif b.kind == "table":
            for r in b.rows[:6]:
                p = tf.add_paragraph()
                p.add_run().text = " | ".join(map(str, r))
    prs.save(str(out))


def _pdf_text(doc: RenderedDoc) -> list[str]:
    return blocks_to_markdown(doc).splitlines()


def render_pdf(doc: RenderedDoc, out: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit
    c = canvas.Canvas(str(out), pagesize=A4)
    width, height = A4
    # 尝试注册 CJK 字体；缺失则回退（ASCII 仍可渲染）。
    try:
        fonts = list(Path("C:/Windows/Fonts").glob("msyh*.ttc"))
        if fonts:
            pdfmetrics.registerFont(TTFont("msyh", str(fonts[0])))
            fname = "msyh"
        else:
            fname = "Helvetica"
    except Exception:
        fname = "Helvetica"
    y = height - 50
    c.setFont(fname, 12)
    c.drawString(50, y, doc.title)
    y -= 22
    c.setFont(fname, 10)
    for line in _pdf_text(doc)[1:1000]:
        for piece in simpleSplit(line, fname, 10, width - 100):
            if y < 40:
                c.showPage()
                c.setFont(fname, 10)
                y = height - 50
            c.drawString(50, y, piece)
            y -= 13
    c.save()


def render(doc: RenderedDoc, out: Path) -> None:
    fmt = doc.format
    if fmt == "docx":
        render_docx(doc, out)
    elif fmt == "xlsx":
        render_xlsx(doc, out)
    elif fmt == "pptx":
        render_pptx(doc, out)
    elif fmt == "pdf":
        render_pdf(doc, out)
    else:
        raise ValueError(f"office renderer got {fmt!r}")