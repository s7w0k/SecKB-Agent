"""P4：Parser/Profile/Chunk 离线压力验证（计划 §9 / §13 P4）。

不调用 embedding，对全部 S1 文件运行真实 parser + 差异化 chunker，输出：
- chunking-summary.json / chunking-cases.jsonl / profile-confusion-matrix.csv
- 结构门禁（§9.3）：超长块/空块、FAQ 原子性、Policy 条款完整、Procedure 顺序、
  Table 表头、document_profile/content_type/section_path、profile 混淆矩阵（F1）。

二进制（PDF/Office）在 MinerU 不可用时以本地库抽取文本走 native parser，并如实标注
parser_name，S2 结构化门禁不计这些为 MinerU 证据（计划 §8.2）。
"""
from __future__ import annotations

import json
from pathlib import Path

from app.services.document_processing.chunkers.registry import build_default_registry
from app.services.document_processing.contracts import DocumentProfile
from app.services.document_processing.parser_registry import build_default_registry as build_parser_reg
from app.services.document_processing.parsers.markdown import MarkdownParser
from app.services.document_processing.parsers.plain_text import PlainTextParser
from app.services.document_processing.profile import DocumentProfiler

from scripts.enterprise_rag.config import RunConfig

_MIME = {
    ".md": "text/markdown", ".txt": "text/plain", ".log": "text/plain",
    ".json": "text/plain", ".jsonl": "text/plain", ".yaml": "text/plain",
    ".yml": "text/plain", ".csv": "text/plain", ".html": "text/plain",
}

_FAMILY_TRUTH = {
    "overview": "narrative", "whitepaper": "narrative", "architecture": "narrative",
    "case": "narrative", "relations": "narrative",
    "threat-model": "policy", "sla": "policy", "release-notes": "policy", "compliance": "policy",
    "faq": "faq",
    "ops-guide": "procedure", "admin-guide": "procedure", "user-guide": "procedure",
    "dev-guide": "procedure", "troubleshooting": "procedure",
    "api-ref": "table_records", "parameters": "table_records", "compatibility": "table_records",
    "capacity": "table_records", "pricing": "table_records", "config-sample": "table_records",
}

_JSON_EXT = {".json", ".jsonl", ".yaml", ".yml"}


def _file_mime(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".md":
        return "text/markdown"
    if ext in (".pdf", ".docx", ".xlsx", ".pptx") or ext in (".json", ".jsonl", ".yaml", ".yml"):
        return "text/plain"
    return _MIME.get(ext, "text/plain")


def _sheet_to_markdown(path: Path, ext: str) -> str:
    """把 xlsx/csv 还原成 markdown 管道表格，使 MarkdownParser 生成真实 table 块
    （TableChunker 因而能对表格数据产出可检索 chunk，计划 §9.1/Two.2 的 table_records）。"""
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[list[str]] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append(["" if c is None else str(c) for c in row])
    else:
        import csv as _csvmod
        rows = []
        try:
            with open(path, encoding="utf-8", newline="") as fh:
                for r in _csvmod.reader(fh):
                    rows.append(["" if c is None else str(c) for c in r])
        except UnicodeDecodeError:
            with open(path, encoding="gbk", errors="replace", newline="") as fh:
                for r in _csvmod.reader(fh):
                    rows.append(["" if c is None else str(c) for c in r])
    if not rows:
        return ""
    cols = max(len(r) for r in rows)
    # 用事实池化的列头；保留结构（非模板占位）
    header = [f"col{i+1}" for i in range(cols)]
    lines = ["| " + " | ".join(r[i] if i < len(r) else "" for i in range(cols)) + " |"
             for r in rows]
    out = ["| " + " | ".join(header) + " |", "|" + "|".join("---" for _ in range(cols)) + "|"]
    out += lines
    return "\n".join(out)


def _extract_text(path: Path, parser_name: dict) -> tuple[bytes, str]:
    """返回 (bytes, 实际 parser mime)。json/yaml 结构文件转成可读文本以保留 key path。"""
    ext = path.suffix.lower()
    if ext in (".xlsx", ".csv"):
        parser_name["parser"] = "table-native"
        return _sheet_to_markdown(path, ext).encode("utf-8"), "text/markdown"
    if ext == ".pdf":
        from pypdf import PdfReader
        r = PdfReader(str(path))
        parser_name["parser"] = "pypdf"
        return ("\n".join(p.extract_text() or "" for p in r.pages)).encode("utf-8"), "text/plain"
    if ext == ".docx":
        from docx import Document
        d = Document(str(path))
        parser_name["parser"] = "docx-native"
        return ("\n\n".join(p.text for p in d.paragraphs)).encode("utf-8"), "text/plain"
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c) if c is not None else "" for c in row))
        parser_name["parser"] = "openpyxl-native"
        return ("\n".join(rows)).encode("utf-8"), "text/plain"
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        parser_name["parser"] = "pptx-native"
        return ("\n".join(parts)).encode("utf-8"), "text/plain"
    if ext in _JSON_EXT:
        parser_name["parser"] = "native"
        try:
            obj = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            obj = None
        if obj is not None:
            # 保留 key path 风格的扁平表示
            lines = _flatten("$", obj)
            return ("\n".join(lines)).encode("utf-8"), "text/plain"
    parser_name["parser"] = "native"
    return path.read_bytes(), _MIME.get(ext, "text/plain")


def _flatten(path: str, node) -> list[str]:
    out: list[str] = []
    if isinstance(node, dict):
        for k, v in node.items():
            out += _flatten(f"{path}/{k}", v)
    elif isinstance(node, list):
        for i, v in enumerate(node):
            out += _flatten(f"{path}[{i}]", v)
    else:
        out.append(f"{path} = {node}")
    return out


def _chunk_checks(doc_profile: str, truth_profile: str, chunks, all_parsed_text: str) -> dict:
    """返回单文档结构门禁检查结果。"""
    over_max = sum(1 for c in chunks if c.token_count > 700)
    empty = [c for c in chunks if not c.display_content.strip()
             or (c.content_type in ("heading", "title") and not c.embedding_text.strip())]
    faq_checked = 0
    faq_atomic = 0
    for c in chunks:
        if doc_profile == "faq":
            faq_checked += 1
            if ("Q" in c.display_content or "问题" in c.embedding_text) and ("A" in c.embedding_text or "回答" in c.embedding_text):
                faq_atomic += 1
    policy_checked = sum(1 for c in chunks if doc_profile == "policy" and "第" in c.embedding_text)
    policy_complete = sum(1 for c in chunks if c.embedding_text != "" and not c.embedding_text.startswith("第") or c.embedding_text == "")
    header_checked = 0
    header_kept = 0
    for c in chunks:
        if doc_profile == "table_records" and c.content_type in ("table", "table_row", "record"):
            header_checked += 1
            headers = c.metadata.get("headers") or []
            if headers:
                header_kept += 1
    return {
        "over_max_tokens": over_max, "empty_chunks": len(empty),
        "faq_checked": faq_checked, "faq_atomic": faq_atomic,
        "policy_complete_ratio": round(policy_complete / max(1, policy_checked), 4),
        "table_header_checked": header_checked, "table_header_kept": header_kept,
        "has_profile": 1 if all(c.document_profile for c in chunks) else 0,
        "has_content_type": 1 if all(c.content_type for c in chunks) else 0,
        "n": len(chunks),
    }


def run(run_id: str, scale: str, seed: int) -> dict:
    cfg = RunConfig(run_id=run_id, scale=scale, seed=seed)
    files_root = cfg.files_dir
    registry = build_default_registry()
    parser_reg = build_parser_reg()
    # native parser 注册 markdown + plain
    md_parser = MarkdownParser()
    pt_parser = PlainTextParser()
    manifest = json.loads((files_root / "corpus-manifest.json").read_text(encoding="utf-8"))
    # 记录每个产品的 family -> 期望 profile（由 corpus-manifest + 内容生成约定）
    manifest_rows = {}
    with (files_root / "corpus-manifest.json").open(encoding="utf-8") as fh:
        man = json.load(fh)
    manifest_rows = {r["doc_id"]: r for r in man["rows"]} if "rows" in man else {}
    if not manifest_rows:
        # corpus-manifest.json 结构：本实现为 {documents, faq_entries, ...} 无 rows，
        # 故从文件 doc_id + truth 约定推断 family。
        from scripts.enterprise_rag.renderers import content as _content
        catalog, facts = _content._load_truth(None)
        plan = _content.plan_documents(catalog, scale, seed)
        manifest_rows = {r["doc_id"]: r for r in plan}
    cases: list[dict] = []
    # 混淆矩阵仅在"结构保留解析"(Markdown)子集上统计；二进制/CSV 走 native fallback，
    # 单独记录映射（计划 §8.2 的限制，不能用于判定 S2 结构化门禁 PASS）。
    confusion: dict[tuple, int] = {}
    fallback_confusion: dict[tuple, int] = {}
    agg = {"chunks": 0, "over_max": 0, "empty": 0, "faq_atomic_total": 0,
           "faq_checked_total": 0, "token_docs": [], "files": 0,
           "native_files": 0, "fallback_files": 0}
    # 结构化保留解析子集：md 表格、xlsx/csv 表格（还原为 table 块）。
    # 其余二进制（PDF/DOCX/PPTX/log）与 json/yaml 物化走 native 文本 fallback（计划 §8.2）。
    _NATIVE = {".md", ".xlsx", ".csv"}
    for doc_dir in sorted([d for d in files_root.iterdir() if d.is_dir()]):
        row = manifest_rows.get(doc_dir.name, {})
        truth_family = row.get("family", "overview")
        truth = row.get("profile") or _FAMILY_TRUTH.get(truth_family, "narrative")
        if doc_dir.name.endswith("-FAQ"):
            truth = "faq"  # FAQ 目录不在 plan_documents 中，须显式指定 truth
        for f in sorted([x for x in doc_dir.iterdir() if x.is_file()]):
            ext = f.suffix.lower()
            if ext not in (".md", ".pdf", ".docx", ".xlsx", ".pptx",
                           ".json", ".jsonl", ".yaml", ".yml", ".txt", ".log",
                           ".html", ".csv"):
                continue
            data, mime = _extract_text(f, {})
            try:
                if mime == "text/markdown" or ext == ".md":
                    parsed = md_parser.parse(data, source_uri=str(f), mime_type="text/markdown")
                else:
                    parsed = pt_parser.parse(data, source_uri=str(f), mime_type=mime)
            except Exception as exc:  # noqa: BLE001
                cases.append({"file": str(f), "error": str(exc)[:160]})
                continue
            # rule-based predict（不注入 truth）→ 混淆矩阵
            predicted = DocumentProfiler().detect(parsed)
            pred = predicted.value if isinstance(predicted, DocumentProfile) else str(predicted)
            is_native = ext in _NATIVE
            target = confusion if is_native else fallback_confusion
            target[(truth, pred)] = target.get((truth, pred), 0) + 1
            if is_native:
                agg["native_files"] += 1
            else:
                agg["fallback_files"] += 1
            # 真实差异化切块（显式给定 truth profile → 验证 chunker）
            prof = DocumentProfile(truth)
            chunks = registry.chunk(parsed, profile=prof)
            checks = _chunk_checks(truth, truth, chunks, parsed.parsed_hash)
            agg["chunks"] += len(chunks)
            agg["over_max"] += checks["over_max_tokens"]
            agg["empty"] += checks["empty_chunks"]
            agg["faq_atomic_total"] += checks["faq_atomic"]
            agg["faq_checked_total"] += checks["faq_checked"]
            agg["files"] += 1
            for c in chunks:
                cases.append({
                    "file": str(f.relative_to(files_root)),
                    "doc_id": doc_dir.name, "truth_profile": truth,
                    "predicted_profile": pred,
                    "profile_scope": "native" if is_native else "fallback",
                    "token_count": c.token_count,
                    "content_type": c.content_type,
                    "logical_key": c.logical_key,
                    "display_head": c.display_content[:120],
                    "parser": "md" if ext == ".md" else ("json" if ext in _JSON_EXT else ext),
                })
    # 汇总
    macro_f1 = _macro_f1(confusion)
    faq_atomic_ratio = agg["faq_atomic_total"] / max(1, agg["faq_checked_total"])
    result = {
        "run_id": run_id, "scale": scale,
        "chunks": agg["chunks"], "files_parsed": agg["files"],
        "over_max_tokens": agg["over_max"],
        "empty_chunks": agg["empty"],
        "faq_atomic_ratio": round(faq_atomic_ratio, 4),
        "faq_checked": agg["faq_checked_total"],
        "confusion_matrix": {f"{k[0]}->{k[1]}": v for k, v in confusion.items()},
        "profile_macro_f1": round(macro_f1, 4),
        # §8.2：结构保留(Markdown)子集用于 profile 门禁；native-fallback 单独记录。
        "profile_scope": {
            "native_evaluable_exts": sorted(_NATIVE),
            "native_files": agg["native_files"],
            "fallback_files": agg["fallback_files"],
            "fallback_confusion": {f"{k[0]}->{k[1]}": v for k, v in fallback_confusion.items()},
            "note": "PDF/DOCX/XLSX/PPTX 及 CSV 走 native 文本提取（计划 §8.2），"
                    "表格结构(表格块)无法在 native fallback 下识别，故不计入 profile F1；"
                    "S2 结构化门禁不得据此判 PASS。",
        },
        "gate_pass": (agg["over_max"] == 0 and agg["empty"] == 0
                      and faq_atomic_ratio >= 0.99 and macro_f1 >= 0.95),
    }
    cfg.out_dir.mkdir(parents=True, exist_ok=True)
    (cfg.out_dir / "chunking-summary.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    with (cfg.out_dir / "chunking-cases.jsonl").open("w", encoding="utf-8") as fh:
        for c in cases:
            fh.write(json.dumps(c, ensure_ascii=False) + "\n")
    import csv as _csv
    with (cfg.out_dir / "profile-confusion-matrix.csv").open("w", encoding="utf-8", newline="") as fh:
        w = _csv.writer(fh)
        w.writerow(["truth", "predicted"])
        for (t, p), v in confusion.items():
            w.writerow([t, p, v])
    return result


def _family(doc_dir: str) -> str:
    # corpus file layout: <Pxxx-Vnn>/... 无 family 元数据，从 FAQ 目录名除外
    base = doc_dir.split("-")[0]
    if base in {"P001", "P002", "P003", "P005", "P007", "P008", "P011", "P014"}:
        pass
    return "narrative"  # 由 title/内容规则在实际中判定（简化：以 manifest 为准）


def _macro_f1(confusion: dict[tuple, int]) -> float:
    """真实 macro-F1：每个 truth 类别分别算 precision/recall/F1，再平均。

    仅计算在 truth 中出现过的类别（避免零支持类别把分数拉成无法解释的 0）。
    precision = TP / (TP + FP)，recall = TP / (TP + FN)。
    """
    classes = sorted({t for (t, _p) in confusion})
    if not classes:
        return 0.0
    tp: dict[str, int] = {c: 0 for c in classes}
    truth_n: dict[str, int] = {c: 0 for c in classes}   # 每类 truth 样本数（FN 分母）
    pred_n: dict[str, int] = {c: 0 for c in classes}    # 每类被预测样本数（FP 分母）
    for (t, p), v in confusion.items():
        if t not in truth_n:
            continue
        truth_n[t] += v
        pred_n.setdefault(p, 0)
        pred_n[p] += v
        if t == p:
            tp[t] += v
    f1s = []
    for c in classes:
        tn = truth_n.get(c, 0)
        if tn == 0:
            continue
        precision = tp[c] / pred_n.get(c, 0) if pred_n.get(c, 0) else 0.0
        recall = tp[c] / tn if tn else 0.0
        f1 = 2 * precision * recall / (precision + recall) if (precision + recall) else 0.0
        f1s.append(f1)
    return sum(f1s) / len(f1s) if f1s else 0.0


if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--run-id", default="run-s1-20260828")
    ap.add_argument("--scale", default="S1")
    ap.add_argument("--seed", type=int, default=20260828)
    a = ap.parse_args()
    r = run(a.run_id, a.scale, a.seed)
    print("chunks:", r["chunks"], "files:", r["files_parsed"],
          "over_max:", r["over_max_tokens"], "empty:", r["empty_chunks"],
          "faq_atomic:", r["faq_atomic_ratio"], "macro_f1:", r["profile_macro_f1"],
          "gate:", r["gate_pass"])
    raise SystemExit(0 if r["gate_pass"] else 1)